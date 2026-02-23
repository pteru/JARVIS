#!/usr/bin/env python3
"""JARVIS PPTX Tool — Create, read, edit, and inspect PowerPoint presentations."""

import argparse
import json
import re
import sys
import unicodedata
import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def resolve_path(filepath):
    """Resolve a file path, handling NFC/NFD Unicode normalization mismatches."""
    p = Path(filepath)
    if p.exists():
        return str(p)
    nfd = unicodedata.normalize("NFD", str(p))
    if os.path.exists(nfd):
        return nfd
    nfc = unicodedata.normalize("NFC", str(p))
    if os.path.exists(nfc):
        return nfc
    return str(p)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

LAYOUT_NAMES = {
    0: "Title Slide",
    1: "Title and Content",
    2: "Section Header",
    3: "Two Content",
    4: "Comparison",
    5: "Title Only",
    6: "Blank",
    7: "Content with Caption",
    8: "Picture with Caption",
}


def shape_to_dict(shape, with_styles=False):
    """Convert a shape to a dict."""
    d = {
        "shape_id": shape.shape_id,
        "name": shape.name,
        "type": shape.shape_type.__str__() if shape.shape_type else "unknown",
        "position": {
            "left": shape.left,
            "top": shape.top,
            "width": shape.width,
            "height": shape.height,
        },
    }

    if shape.has_text_frame:
        d["text"] = shape.text_frame.text
        if with_styles:
            d["paragraphs"] = []
            for para in shape.text_frame.paragraphs:
                para_dict = {
                    "text": para.text,
                    "level": para.level,
                    "runs": [],
                }
                if para.alignment is not None:
                    para_dict["alignment"] = str(para.alignment)
                for run in para.runs:
                    run_dict = {"text": run.text}
                    if run.font.bold:
                        run_dict["bold"] = True
                    if run.font.italic:
                        run_dict["italic"] = True
                    if run.font.underline:
                        run_dict["underline"] = True
                    if run.font.size:
                        run_dict["size"] = run.font.size.pt
                    if run.font.name:
                        run_dict["font"] = run.font.name
                    if run.font.color and run.font.color.rgb:
                        run_dict["color"] = str(run.font.color.rgb)
                    para_dict["runs"].append(run_dict)
                d["paragraphs"].append(para_dict)

    if shape.has_table:
        table = shape.table
        rows = []
        for row in table.rows:
            rows.append([cell.text for cell in row.cells])
        d["table"] = {"rows": rows, "row_count": len(table.rows), "col_count": len(table.columns)}

    if hasattr(shape, "image"):
        try:
            img = shape.image
            d["image"] = {
                "content_type": img.content_type,
                "size": len(img.blob),
            }
        except Exception:
            pass

    return d


def slide_to_dict(slide, idx, with_styles=False):
    """Convert a slide to a dict."""
    layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"
    d = {
        "slide_number": idx + 1,
        "layout": layout_name,
        "shapes": [shape_to_dict(s, with_styles) for s in slide.shapes],
    }

    # Extract notes
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        notes_text = slide.notes_slide.notes_text_frame.text
        if notes_text.strip():
            d["notes"] = notes_text

    return d


def get_metadata(prs):
    """Extract core properties as a dict."""
    cp = prs.core_properties
    meta = {}
    for attr in ("author", "title", "subject", "keywords", "category",
                 "comments", "last_modified_by", "revision"):
        val = getattr(cp, attr, None)
        if val:
            meta[attr] = str(val)
    for attr in ("created", "modified"):
        val = getattr(cp, attr, None)
        if val:
            meta[attr] = val.isoformat()
    return meta


# ---------------------------------------------------------------------------
# Markdown conversion
# ---------------------------------------------------------------------------

def _runs_to_markdown(runs):
    """Convert runs with inline formatting to markdown text."""
    parts = []
    for run in runs:
        text = run.text
        if not text:
            continue
        if run.font.bold and run.font.italic:
            text = f"***{text}***"
        elif run.font.bold:
            text = f"**{text}**"
        elif run.font.italic:
            text = f"*{text}*"
        parts.append(text)
    return "".join(parts)


def slide_to_markdown(slide, idx):
    """Convert a slide to markdown."""
    lines = []
    layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"
    lines.append(f"## Slide {idx + 1} — {layout_name}")
    lines.append("")

    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = _runs_to_markdown(para.runs)
                if not text.strip():
                    continue
                # Title placeholder → heading
                if shape.is_placeholder and shape.placeholder_format.idx in (0, 1):
                    if shape.placeholder_format.idx == 0:
                        lines.append(f"### {text}")
                    else:
                        lines.append(text)
                elif para.level > 0:
                    indent = "  " * para.level
                    lines.append(f"{indent}- {text}")
                else:
                    lines.append(text)
            lines.append("")

        if shape.has_table:
            table = shape.table
            rows = []
            for row in table.rows:
                rows.append([cell.text.replace("\n", " ") for cell in row.cells])
            if rows:
                lines.append("| " + " | ".join(rows[0]) + " |")
                lines.append("| " + " | ".join(["---"] * len(rows[0])) + " |")
                for row in rows[1:]:
                    padded = row + [""] * (len(rows[0]) - len(row))
                    lines.append("| " + " | ".join(padded) + " |")
                lines.append("")

    # Notes
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            lines.append(f"> **Notes:** {notes}")
            lines.append("")

    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Commands
# ---------------------------------------------------------------------------

def cmd_read(args):
    prs = Presentation(resolve_path(args.file))

    slides = []
    for idx, slide in enumerate(prs.slides):
        slides.append(slide_to_dict(slide, idx, args.with_styles))

    # Apply range filter
    if args.range:
        parts = args.range.split(":")
        start = int(parts[0]) - 1 if parts[0] else 0  # 1-based to 0-based
        end = int(parts[1]) if len(parts) > 1 and parts[1] else len(slides)
        slides = slides[start:end]

    if args.format == "json":
        result = {
            "metadata": get_metadata(prs),
            "slide_count": len(prs.slides),
            "slide_width": prs.slide_width,
            "slide_height": prs.slide_height,
            "slides": slides,
        }
        print(json.dumps(result, default=str, indent=2, ensure_ascii=False))

    elif args.format == "markdown":
        prs2 = Presentation(resolve_path(args.file))
        slide_list = list(prs2.slides)
        if args.range:
            parts = args.range.split(":")
            start = int(parts[0]) - 1 if parts[0] else 0
            end = int(parts[1]) if len(parts) > 1 and parts[1] else len(slide_list)
            slide_list = slide_list[start:end]
            start_idx = start
        else:
            start_idx = 0

        md_lines = []
        meta = get_metadata(prs2)
        if meta.get("title"):
            md_lines.append(f"# {meta['title']}")
            md_lines.append("")

        for i, slide in enumerate(slide_list):
            md_lines.append(slide_to_markdown(slide, start_idx + i))

        print("\n".join(md_lines))

    else:
        # Plain text
        for s in slides:
            print(f"--- Slide {s['slide_number']} ({s['layout']}) ---")
            for shape in s["shapes"]:
                if "text" in shape and shape["text"].strip():
                    print(shape["text"])
                if "table" in shape:
                    for row in shape["table"]["rows"]:
                        print("\t".join(row))
            if "notes" in s:
                print(f"[Notes: {s['notes']}]")
            print()


def cmd_create(args):
    # Determine input source
    if args.input:
        input_path = Path(args.input)
        content = input_path.read_text(encoding="utf-8")
        is_json = input_path.suffix.lower() == ".json"
    else:
        content = sys.stdin.read()
        stripped = content.strip()
        is_json = stripped.startswith("{") or stripped.startswith("[")

    if args.template:
        prs = Presentation(resolve_path(args.template))
    else:
        prs = Presentation()

    if is_json:
        _create_from_json(prs, json.loads(content))
    else:
        _create_from_markdown(prs, content)

    prs.save(args.file)
    print(f"Created {args.file}")


def _get_layout(prs, layout_name_or_idx):
    """Get a slide layout by name or index."""
    if isinstance(layout_name_or_idx, int):
        try:
            return prs.slide_layouts[layout_name_or_idx]
        except IndexError:
            return prs.slide_layouts[6]  # Blank fallback

    name_lower = layout_name_or_idx.lower()
    for layout in prs.slide_layouts:
        if layout.name.lower() == name_lower:
            return layout

    # Try common aliases
    aliases = {
        "title": 0, "title_slide": 0,
        "title_and_content": 1, "content": 1,
        "section": 2, "section_header": 2,
        "two_content": 3,
        "comparison": 4,
        "title_only": 5,
        "blank": 6,
    }
    idx = aliases.get(name_lower.replace(" ", "_"))
    if idx is not None:
        try:
            return prs.slide_layouts[idx]
        except IndexError:
            pass

    return prs.slide_layouts[1]  # Default: Title and Content


def _create_from_json(prs, data):
    """Build presentation from JSON structure."""
    if isinstance(data, dict):
        if "metadata" in data:
            meta = data["metadata"]
            cp = prs.core_properties
            for key in ("author", "title", "subject", "keywords",
                        "category", "comments"):
                if key in meta:
                    setattr(cp, key, meta[key])
        slides_data = data.get("slides", [])
    elif isinstance(data, list):
        slides_data = data
    else:
        return

    for slide_data in slides_data:
        if not isinstance(slide_data, dict):
            continue

        layout_ref = slide_data.get("layout", 1)
        layout = _get_layout(prs, layout_ref)
        slide = prs.slides.add_slide(layout)

        # Set title
        if "title" in slide_data and slide.shapes.title:
            slide.shapes.title.text = slide_data["title"]

        # Set body / content placeholder
        if "body" in slide_data or "content" in slide_data:
            body_text = slide_data.get("body", slide_data.get("content", ""))
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:  # Content placeholder
                    ph.text = body_text
                    break

        # Bullet points
        if "bullets" in slide_data:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    tf = ph.text_frame
                    tf.clear()
                    for i, bullet in enumerate(slide_data["bullets"]):
                        if isinstance(bullet, dict):
                            text = bullet.get("text", "")
                            level = bullet.get("level", 0)
                        else:
                            text = str(bullet)
                            level = 0
                        if i == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        p.text = text
                        p.level = level
                    break

        # Table
        if "table" in slide_data:
            tdata = slide_data["table"]
            rows_data = tdata.get("rows", tdata) if isinstance(tdata, dict) else tdata
            if rows_data:
                nrows = len(rows_data)
                ncols = max(len(r) for r in rows_data)
                left = Inches(0.5)
                top = Inches(2.0)
                width = Inches(9.0)
                height = Inches(0.8 * nrows)
                table = slide.shapes.add_table(nrows, ncols, left, top, width, height).table
                for r, row in enumerate(rows_data):
                    for c, val in enumerate(row):
                        if c < ncols:
                            table.cell(r, c).text = str(val)

        # Notes
        if "notes" in slide_data:
            slide.notes_slide.notes_text_frame.text = slide_data["notes"]


def _create_from_markdown(prs, text):
    """Build presentation from markdown text. Each ## heading starts a new slide."""
    lines = text.split("\n")
    slide_chunks = []
    current = None

    for line in lines:
        # Top-level heading → presentation title slide
        m = re.match(r"^#\s+(.*)", line)
        if m and not re.match(r"^##\s+", line):
            if current:
                slide_chunks.append(current)
            current = {"title": m.group(1).strip(), "layout": 0, "lines": []}
            continue

        # Slide heading
        m = re.match(r"^##\s+(.*)", line)
        if m:
            if current:
                slide_chunks.append(current)
            current = {"title": m.group(1).strip(), "layout": 1, "lines": []}
            continue

        if current is None:
            current = {"title": "", "layout": 6, "lines": []}

        current["lines"].append(line)

    if current:
        slide_chunks.append(current)

    for chunk in slide_chunks:
        layout = _get_layout(prs, chunk["layout"])
        slide = prs.slides.add_slide(layout)

        if chunk["title"] and slide.shapes.title:
            slide.shapes.title.text = chunk["title"]

        # Process body lines
        body_lines = chunk["lines"]
        # Strip leading/trailing empty lines
        while body_lines and not body_lines[0].strip():
            body_lines.pop(0)
        while body_lines and not body_lines[-1].strip():
            body_lines.pop()

        if not body_lines:
            continue

        # Check if it's a table
        table_lines = [l for l in body_lines if l.strip().startswith("|")]
        if len(table_lines) > 1 and len(table_lines) == len([l for l in body_lines if l.strip()]):
            _add_table_to_slide(slide, table_lines)
            continue

        # Find content placeholder
        content_ph = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                content_ph = ph
                break

        if content_ph is None:
            # No content placeholder — add a text box
            left = Inches(0.5)
            top = Inches(1.8)
            width = Inches(9.0)
            height = Inches(5.0)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            content_ph = txBox

        tf = content_ph.text_frame
        tf.clear()

        first = True
        for line in body_lines:
            if not line.strip():
                continue

            # Bullet
            bm = re.match(r"^(\s*)[-*]\s+(.*)", line)
            # Numbered
            nm = re.match(r"^(\s*)\d+\.\s+(.*)", line)
            # Sub-heading
            hm = re.match(r"^###\s+(.*)", line)
            # Blockquote (notes)
            qm = re.match(r"^>\s+(.*)", line)

            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()

            if hm:
                p.text = hm.group(1).strip()
                if p.runs:
                    p.runs[0].font.bold = True
                else:
                    run = p.add_run()
                    run.text = hm.group(1).strip()
                    run.font.bold = True
                    # Clear plain text
                    p._p.text = ""
                    p._p.append(run._r)
            elif bm:
                indent_level = len(bm.group(1)) // 2
                _add_formatted_runs(p, bm.group(2).strip())
                p.level = indent_level
            elif nm:
                indent_level = len(nm.group(1)) // 2
                _add_formatted_runs(p, nm.group(2).strip())
                p.level = indent_level
            elif qm:
                # Treat as notes on the slide
                slide.notes_slide.notes_text_frame.text = qm.group(1).strip()
            else:
                _add_formatted_runs(p, line.strip())


def _add_table_to_slide(slide, lines):
    """Parse markdown table lines and add to slide."""
    rows = []
    for line in lines:
        cells = [c.strip() for c in line.strip().strip("|").split("|")]
        if all(re.match(r"^[-:]+$", c) for c in cells):
            continue
        rows.append(cells)
    if not rows:
        return
    ncols = max(len(r) for r in rows)
    nrows = len(rows)
    left = Inches(0.5)
    top = Inches(2.0)
    width = Inches(9.0)
    height = Inches(0.8 * nrows)
    table = slide.shapes.add_table(nrows, ncols, left, top, width, height).table
    for r, row in enumerate(rows):
        for c, text in enumerate(row):
            if c < ncols:
                table.cell(r, c).text = text


def _add_formatted_runs(paragraph, text):
    """Parse inline markdown formatting into runs on a pptx paragraph."""
    paragraph.clear()
    pattern = re.compile(r"(\*{3}(.+?)\*{3}|\*{2}(.+?)\*{2}|\*(.+?)\*)")
    pos = 0
    for m in pattern.finditer(text):
        if m.start() > pos:
            run = paragraph.add_run()
            run.text = text[pos:m.start()]
        if m.group(2):  # ***bold italic***
            run = paragraph.add_run()
            run.text = m.group(2)
            run.font.bold = True
            run.font.italic = True
        elif m.group(3):  # **bold**
            run = paragraph.add_run()
            run.text = m.group(3)
            run.font.bold = True
        elif m.group(4):  # *italic*
            run = paragraph.add_run()
            run.text = m.group(4)
            run.font.italic = True
        pos = m.end()
    if pos < len(text):
        run = paragraph.add_run()
        run.text = text[pos:]


def cmd_edit(args):
    prs = Presentation(resolve_path(args.file))

    if args.replace or args.replace_all:
        old_text, new_text = (args.replace or args.replace_all)
        replace_all = args.replace_all is not None
        count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if old_text in run.text:
                                run.text = run.text.replace(old_text, new_text, 1 if not replace_all else -1)
                                count += 1
                                if not replace_all:
                                    break
                        if count > 0 and not replace_all:
                            break
                    if count > 0 and not replace_all:
                        break
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text_frame:
                                for para in cell.text_frame.paragraphs:
                                    for run in para.runs:
                                        if old_text in run.text:
                                            run.text = run.text.replace(old_text, new_text, 1 if not replace_all else -1)
                                            count += 1
                                            if not replace_all:
                                                break
                                    if count > 0 and not replace_all:
                                        break
                                if count > 0 and not replace_all:
                                    break
                        if count > 0 and not replace_all:
                            break
            if count > 0 and not replace_all:
                break
        print(f"Replaced {count} occurrence(s)")

    if args.add_slide is not None:
        layout_ref = args.add_slide
        try:
            layout_ref = int(layout_ref)
        except ValueError:
            pass
        layout = _get_layout(prs, layout_ref)
        slide = prs.slides.add_slide(layout)
        title_text = args.title or ""
        if title_text and slide.shapes.title:
            slide.shapes.title.text = title_text
        body_text = args.body or ""
        if body_text:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    ph.text = body_text
                    break
        print(f"Added slide (layout: {layout.name})")

    if args.delete_slide is not None:
        idx = args.delete_slide - 1  # 1-based to 0-based
        slides = list(prs.slides)
        if 0 <= idx < len(slides):
            rId = prs.slides._sldIdLst[idx].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[idx]
            print(f"Deleted slide {args.delete_slide}")
        else:
            print(f"Error: slide {args.delete_slide} out of range (1-{len(slides)})", file=sys.stderr)
            sys.exit(1)

    if args.set_notes is not None:
        slide_num, notes_text = args.set_notes
        slide_num = int(slide_num)
        slides = list(prs.slides)
        if 1 <= slide_num <= len(slides):
            slide = slides[slide_num - 1]
            slide.notes_slide.notes_text_frame.text = notes_text
            print(f"Set notes on slide {slide_num}")
        else:
            print(f"Error: slide {slide_num} out of range", file=sys.stderr)
            sys.exit(1)

    if args.set_metadata:
        key, value = args.set_metadata
        cp = prs.core_properties
        if hasattr(cp, key):
            setattr(cp, key, value)
            print(f"Set metadata '{key}' = '{value}'")
        else:
            print(f"Error: unknown metadata key '{key}'", file=sys.stderr)
            sys.exit(1)

    prs.save(args.file)
    print(f"Saved {args.file}")


def cmd_info(args):
    prs = Presentation(resolve_path(args.file))
    meta = get_metadata(prs)

    slide_count = len(prs.slides)
    layouts_used = set()
    total_shapes = 0
    total_tables = 0
    total_images = 0
    word_count = 0

    for slide in prs.slides:
        if slide.slide_layout:
            layouts_used.add(slide.slide_layout.name)
        for shape in slide.shapes:
            total_shapes += 1
            if shape.has_table:
                total_tables += 1
            if hasattr(shape, "image"):
                try:
                    _ = shape.image
                    total_images += 1
                except Exception:
                    pass
            if shape.has_text_frame:
                word_count += len(shape.text_frame.text.split())

    info = {
        "file": args.file,
        "metadata": meta,
        "slide_count": slide_count,
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "total_shapes": total_shapes,
        "tables": total_tables,
        "images": total_images,
        "word_count": word_count,
        "layouts_used": sorted(layouts_used),
    }

    if args.format == "json":
        print(json.dumps(info, default=str, indent=2, ensure_ascii=False))
    else:
        print(f"File: {args.file}")
        if meta.get("title"):
            print(f"Title: {meta['title']}")
        if meta.get("author"):
            print(f"Author: {meta['author']}")
        if meta.get("created"):
            print(f"Created: {meta['created']}")
        if meta.get("modified"):
            print(f"Modified: {meta['modified']}")
        print(f"Slides: {slide_count}")
        print(f"Shapes: {total_shapes}")
        print(f"Tables: {total_tables}")
        print(f"Images: {total_images}")
        print(f"Word count: ~{word_count}")
        dims = f"{prs.slide_width / 914400:.1f}\" x {prs.slide_height / 914400:.1f}\""
        print(f"Dimensions: {dims}")
        if layouts_used:
            print(f"Layouts: {', '.join(sorted(layouts_used))}")


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(description="JARVIS PPTX Tool")
    sub = parser.add_subparsers(dest="command", required=True)

    # read
    p = sub.add_parser("read", help="Read PPTX content")
    p.add_argument("file")
    p.add_argument("--format", choices=["text", "json", "markdown"],
                   default="text")
    p.add_argument("--with-styles", action="store_true",
                   help="Include font, color, bold/italic metadata (JSON only)")
    p.add_argument("--range", help="Slide range e.g. 1:5 (1-based)")

    # create
    p = sub.add_parser("create", help="Create PPTX from input")
    p.add_argument("file")
    p.add_argument("--input", "-i", help="Input file (markdown or JSON)")
    p.add_argument("--template", help="Template PPTX to use as base")

    # edit
    p = sub.add_parser("edit", help="Edit PPTX file")
    p.add_argument("file")
    p.add_argument("--replace", nargs=2, metavar=("OLD", "NEW"),
                   help="Replace first occurrence of text")
    p.add_argument("--replace-all", nargs=2, metavar=("OLD", "NEW"),
                   help="Replace all occurrences of text")
    p.add_argument("--add-slide", metavar="LAYOUT",
                   help="Add slide with layout (name or index)")
    p.add_argument("--title", help="Title for added slide")
    p.add_argument("--body", help="Body text for added slide")
    p.add_argument("--delete-slide", type=int, metavar="NUM",
                   help="Delete slide by number (1-based)")
    p.add_argument("--set-notes", nargs=2, metavar=("SLIDE_NUM", "TEXT"),
                   help="Set speaker notes on a slide")
    p.add_argument("--set-metadata", nargs=2, metavar=("KEY", "VALUE"),
                   help="Set presentation metadata field")

    # info
    p = sub.add_parser("info", help="Show PPTX info")
    p.add_argument("file")
    p.add_argument("--format", choices=["text", "json"], default="text")

    args = parser.parse_args()
    {"read": cmd_read, "create": cmd_create, "edit": cmd_edit,
     "info": cmd_info}[args.command](args)


if __name__ == "__main__":
    main()
